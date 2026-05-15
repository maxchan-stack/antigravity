import os
import shutil
from typing import List
from fastapi import UploadFile

class FileService:
    @staticmethod
    def read_file(path: str) -> str:
        """
        Reads the content of a file based on its extension.
        Supports: .txt, .md, .py, .json, .csv, .docx, .pdf, .xlsx, .pptx
        """
        ext = os.path.splitext(path)[1].lower()
        
        if ext == '.docx':
            from docx import Document
            doc = Document(path)
            return '\n'.join([para.text for para in doc.paragraphs])
        
        if ext == '.pdf':
            import pdfplumber
            text_content = []
            try:
                with pdfplumber.open(path) as pdf:
                    for page in pdf.pages:
                        extracted = page.extract_text()
                        if extracted:
                            text_content.append(extracted)
            except Exception as e:
                # Fallback to pypdf if pdfplumber fails (e.g. PDFObjRef error)
                print(f"pdfplumber failed: {e}. Falling back to pypdf.")
                from pypdf import PdfReader
                reader = PdfReader(path)
                text_content = []
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text_content.append(extracted)
            
            return '\n'.join(text_content)

        if ext in ['.xlsx', '.xls', '.xlsm']:
            import pandas as pd
            # Read all sheets
            xls = pd.ExcelFile(path)
            text_parts = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(df.to_csv(index=False))
            return '\n'.join(text_parts)
        
        if ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(path)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            return '\n'.join(text_content)

        if ext == '.ppt':
            raise ValueError("Legacy .ppt files are not supported. Please save as .pptx and upload again.")

        # Default text read
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    @staticmethod
    def cleanup_files(file_paths: List[str]):
        """
        Deletes the specified files from the filesystem.
        To be used as a BackgroundTask.
        """
        for path in file_paths:
            try:
                if os.path.exists(path):
                    os.remove(path)
                    print(f"Cleaned up file: {path}")
            except Exception as e:
                print(f"Error cleaning up file {path}: {e}")

    @staticmethod
    def read_excel_struct(file_path: str) -> dict:
        """
        Reads Excel file into a dictionary where keys are sheet names 
        and values are lists of lists (rows).
        Values are converted to strings for consistent comparison.
        """
        import pandas as pd
        import numpy as np
        
        # Read all sheets, keep NaN as is to handle later or fillna
        # dtype=str to avoid float conversion issues with IDs etc.
        try:
            df_dict = pd.read_excel(file_path, sheet_name=None, dtype=str)
        except Exception as e:
            print(f"Error reading Excel struct: {e}")
            return {}
        
        result = {}
        for sheet, df in df_dict.items():
            # Replace nan with empty string
            df = df.replace({np.nan: ""})
            # Convert to list of lists (including header)
            header = df.columns.tolist()
            rows = df.values.tolist()
            result[sheet] = [header] + rows
            
        return result
